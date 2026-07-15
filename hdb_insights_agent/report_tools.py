import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Define design constants for a premium aesthetic
COLOR_DARK_SLATE = RGBColor(26, 37, 48)     # #1A2530
COLOR_GOLD = RGBColor(212, 175, 55)         # #D4AF37
COLOR_OFF_WHITE = RGBColor(248, 249, 250)   # #F8F9FA
COLOR_CHARCOAL = RGBColor(60, 60, 60)       # #3C3C3C
COLOR_WHITE = RGBColor(255, 255, 255)

FONT_TITLE = "Trebuchet MS"
FONT_BODY = "Calibri"

def set_slide_background(slide, color):
    """Set solid color background for a slide."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_title_slide(prs, title, subtitle, notes=None):
    """Add a premium title slide with a dark slate background."""
    blank_layout = prs.slide_layouts[6] # Blank slide
    slide = prs.slides.add_slide(blank_layout)
    set_slide_background(slide, COLOR_DARK_SLATE)
    
    # Title & Subtitle in a single text frame to avoid overlaps
    left = Inches(1.0)
    top = Inches(2.2)
    width = Inches(11.3)
    height = Inches(4.0)
    
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0)
    tf.margin_top = Inches(0)
    
    # Main Title
    p_title = tf.paragraphs[0]
    p_title.text = title
    p_title.font.name = FONT_TITLE
    p_title.font.size = Pt(44)
    p_title.font.bold = True
    p_title.font.color.rgb = COLOR_WHITE
    p_title.space_after = Pt(20)
    
    # Subtitle
    if subtitle:
        p_sub = tf.add_paragraph()
        p_sub.text = subtitle
        p_sub.font.name = FONT_BODY
        p_sub.font.size = Pt(20)
        p_sub.font.color.rgb = COLOR_GOLD
        
    if notes:
        slide.notes_slide.notes_text_frame.text = notes
        
    return slide

def add_content_slide(prs, title, bullets, notes=None):
    """Add a clean light-themed content slide with bullet points."""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    set_slide_background(slide, COLOR_OFF_WHITE)
    
    # Slide Title
    left_title = Inches(0.8)
    top_title = Inches(0.6)
    width_title = Inches(11.7)
    height_title = Inches(1.0)
    
    txBox_title = slide.shapes.add_textbox(left_title, top_title, width_title, height_title)
    tf_title = txBox_title.text_frame
    tf_title.word_wrap = True
    tf_title.margin_left = Inches(0)
    tf_title.margin_top = Inches(0)
    
    p_title = tf_title.paragraphs[0]
    p_title.text = title
    p_title.font.name = FONT_TITLE
    p_title.font.size = Pt(36)
    p_title.font.bold = True
    p_title.font.color.rgb = COLOR_DARK_SLATE
    
    # Optional gold decorative line under the title
    # We can use slide margin spacing instead, which is simpler and cleaner.
    
    # Content Bullets
    if bullets:
        left_body = Inches(0.8)
        top_body = Inches(1.8)
        width_body = Inches(11.7)
        height_body = Inches(4.8)
        
        txBox_body = slide.shapes.add_textbox(left_body, top_body, width_body, height_body)
        tf_body = txBox_body.text_frame
        tf_body.word_wrap = True
        tf_body.margin_left = Inches(0)
        tf_body.margin_top = Inches(0)
        
        for i, bullet in enumerate(bullets):
            p = tf_body.paragraphs[0] if i == 0 else tf_body.add_paragraph()
            p.text = bullet
            p.level = 0
            p.font.name = FONT_BODY
            p.font.size = Pt(16)
            p.font.color.rgb = COLOR_CHARCOAL
            p.space_after = Pt(12)
            p.line_spacing = 1.2
            
    if notes:
        slide.notes_slide.notes_text_frame.text = notes
        
    return slide

def add_two_column_slide(prs, title, col1_title, col1_bullets, col2_title, col2_bullets, notes=None):
    """Add a content slide split into two columns."""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    set_slide_background(slide, COLOR_OFF_WHITE)
    
    # Slide Title
    left_title = Inches(0.8)
    top_title = Inches(0.6)
    width_title = Inches(11.7)
    height_title = Inches(1.0)
    
    txBox_title = slide.shapes.add_textbox(left_title, top_title, width_title, height_title)
    tf_title = txBox_title.text_frame
    tf_title.word_wrap = True
    tf_title.margin_left = Inches(0)
    tf_title.margin_top = Inches(0)
    
    p_title = tf_title.paragraphs[0]
    p_title.text = title
    p_title.font.name = FONT_TITLE
    p_title.font.size = Pt(36)
    p_title.font.bold = True
    p_title.font.color.rgb = COLOR_DARK_SLATE
    
    # Column 1
    left_col1 = Inches(0.8)
    top_col1 = Inches(1.8)
    width_col1 = Inches(5.6)
    height_col1 = Inches(4.8)
    
    txBox_col1 = slide.shapes.add_textbox(left_col1, top_col1, width_col1, height_col1)
    tf_col1 = txBox_col1.text_frame
    tf_col1.word_wrap = True
    tf_col1.margin_left = Inches(0)
    tf_col1.margin_top = Inches(0)
    
    if col1_title:
        p_c1_t = tf_col1.paragraphs[0]
        p_c1_t.text = col1_title
        p_c1_t.font.name = FONT_TITLE
        p_c1_t.font.size = Pt(20)
        p_c1_t.font.bold = True
        p_c1_t.font.color.rgb = COLOR_GOLD
        p_c1_t.space_after = Pt(10)
        
    if col1_bullets:
        start_idx = 1 if col1_title else 0
        for i, bullet in enumerate(col1_bullets):
            p = tf_col1.add_paragraph() if (i > 0 or start_idx == 1) else tf_col1.paragraphs[0]
            p.text = bullet
            p.level = 0
            p.font.name = FONT_BODY
            p.font.size = Pt(15)
            p.font.color.rgb = COLOR_CHARCOAL
            p.space_after = Pt(8)
            p.line_spacing = 1.15

    # Column 2
    left_col2 = Inches(6.8)
    top_col2 = Inches(1.8)
    width_col2 = Inches(5.6)
    height_col2 = Inches(4.8)
    
    txBox_col2 = slide.shapes.add_textbox(left_col2, top_col2, width_col2, height_col2)
    tf_col2 = txBox_col2.text_frame
    tf_col2.word_wrap = True
    tf_col2.margin_left = Inches(0)
    tf_col2.margin_top = Inches(0)
    
    if col2_title:
        p_c2_t = tf_col2.paragraphs[0]
        p_c2_t.text = col2_title
        p_c2_t.font.name = FONT_TITLE
        p_c2_t.font.size = Pt(20)
        p_c2_t.font.bold = True
        p_c2_t.font.color.rgb = COLOR_GOLD
        p_c2_t.space_after = Pt(10)
        
    if col2_bullets:
        start_idx = 1 if col2_title else 0
        for i, bullet in enumerate(col2_bullets):
            p = tf_col2.add_paragraph() if (i > 0 or start_idx == 1) else tf_col2.paragraphs[0]
            p.text = bullet
            p.level = 0
            p.font.name = FONT_BODY
            p.font.size = Pt(15)
            p.font.color.rgb = COLOR_CHARCOAL
            p.space_after = Pt(8)
            p.line_spacing = 1.15
            
    if notes:
        slide.notes_slide.notes_text_frame.text = notes
        
    return slide

def generate_powerpoint_report(slides: list[dict], filename: str = "business_insights_report.pptx") -> str:
    """
    Generate a beautiful widescreen PowerPoint (.pptx) presentation from the given slides.
    
    The output is saved inside the 'reports' folder in the project workspace.
    
    Args:
        slides: A list of slide dictionaries, where each slide can have keys:
            - 'slide_type': str ('title', 'content', 'two_column', or 'conclusion')
            - 'title': str
            - 'subtitle': str (optional, for title/conclusion)
            - 'bullets': list[str] (optional, for content/conclusion)
            - 'col1_title': str (optional, for two_column)
            - 'col1_bullets': list[str] (optional, for two_column)
            - 'col2_title': str (optional, for two_column)
            - 'col2_bullets': list[str] (optional, for two_column)
            - 'notes': str (optional speaker notes)
        filename: Optional custom filename ending in '.pptx' (defaults to 'business_insights_report.pptx')
        
    Returns:
        A message string confirming the save path of the file.
    """
    if not filename.endswith('.pptx'):
        filename += '.pptx'
        
    # Ensure folder 'reports' exists in workspace
    workspace_dir = os.path.dirname(os.path.abspath(__file__))
    reports_dir = os.path.join(workspace_dir, "reports")
    os.makedirs(reports_dir, exist_ok=True)
    
    file_path = os.path.join(reports_dir, filename)
    
    # Initialize Presentation
    prs = Presentation()
    
    # Widescreen 16:9
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    for slide_data in slides:
        slide_type = slide_data.get('slide_type', 'content').lower()
        title = slide_data.get('title', 'Untitled Slide')
        notes = slide_data.get('notes')
        
        if slide_type == 'title':
            subtitle = slide_data.get('subtitle', '')
            add_title_slide(prs, title, subtitle, notes)
        elif slide_type == 'two_column':
            col1_title = slide_data.get('col1_title', '')
            col1_bullets = slide_data.get('col1_bullets', [])
            col2_title = slide_data.get('col2_title', '')
            col2_bullets = slide_data.get('col2_bullets', [])
            add_two_column_slide(prs, title, col1_title, col1_bullets, col2_title, col2_bullets, notes)
        elif slide_type == 'conclusion':
            bullets = slide_data.get('bullets', [])
            add_content_slide(prs, title, bullets, notes)
            # Conclusion slides can use content layout or custom styling
        else: # Default is content slide
            bullets = slide_data.get('bullets', [])
            add_content_slide(prs, title, bullets, notes)
            
    prs.save(file_path)
    
    return f"Successfully generated PowerPoint report at: {file_path}"
